import os
import shutil
from datetime import datetime
from pptx import Presentation
from dotenv import load_dotenv

load_dotenv()


CARPETA_ORIGEN = os.getenv('CARPETA_INICIAL')
CARPETA_DESTINO = os.getenv('CARPETA_FINAL')
def buscar_etiqueta(cod_articulo):
    """Busca el archivo de etiqueta correspondiente al código de artículo"""
    archivo_original = os.path.join(CARPETA_ORIGEN, f"{cod_articulo}.pptx")
    print(archivo_original)
    if os.path.exists(archivo_original):
        return archivo_original
    raise Exception(f"No se encontró el archivo {archivo_original}")

def modificar_etiqueta(archivo, nro_lote, fecha_vencimiento):
    """Modifica el contenido de la etiqueta"""
    try:
        # Verificar que el archivo existe
        if not os.path.exists(archivo):
            raise Exception(f"El archivo {archivo} no existe")

        prs = Presentation(archivo)
        
        # Buscar y reemplazar los textos en todas las diapositivas
        for slide in prs.slides:
            # Buscar en todas las formas
            for shape in slide.shapes:
                # Buscar en el texto de la forma
                if hasattr(shape, "text"):
                    texto = shape.text.strip()
                    print(f"Texto encontrado en shape: '{texto}'")
                    if "XXXXXXX" in texto:
                        # Preservar el formato del texto original
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if "XXXXXXX" in run.text:
                                        run.text = run.text.replace("XXXXXXX", str(nro_lote))

                    if "XX-XX-XXXX" in texto:
                        # Preservar el formato del texto original
                        if hasattr(shape, "text_frame"):
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if "XX-XX-XXXX" in run.text:
                                        run.text = run.text.replace("XX-XX-XXXX", fecha_vencimiento)
                
            
                
                # Buscar en grupos de formas
                if shape.shape_type == 6:  
                    for sub_shape in shape.shapes:
                        if hasattr(sub_shape, "text"):
                            texto = sub_shape.text.strip()
                            if "XXXXXXX" in texto:
                                # Preservar el formato del texto en el grupo
                                if hasattr(sub_shape, "text_frame"):
                                    for paragraph in sub_shape.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            if "XXXXXXX" in run.text:
                                                run.text = run.text.replace("XXXXXXX", str(nro_lote))
                            if "XX-XX-XXXX" in texto:
                                # Preservar el formato del texto en el grupo
                                if hasattr(sub_shape, "text_frame"):
                                    for paragraph in sub_shape.text_frame.paragraphs:
                                        for run in paragraph.runs:
                                            if "XX-XX-XXXX" in run.text:
                                                run.text = run.text.replace("XX-XX-XXXX", fecha_vencimiento)

        
        prs.save(archivo)

    except Exception as e:
        raise Exception(f"Error al modificar la etiqueta: {str(e)}")

def procesar_etiqueta(cod_articulo, nro_lote, fecha_vencimiento):
    """Procesa la etiqueta para su impresión"""
    try:
        # Buscar el archivo de etiqueta original
        archivo_original = buscar_etiqueta(cod_articulo)
        
        # Construir la ruta del archivo de destino
        archivo_destino = os.path.join(CARPETA_DESTINO, f"{cod_articulo}.pptx")

        # Asegurarse de que la carpeta de destino existe
        os.makedirs(CARPETA_DESTINO, exist_ok=True)

        # Copiar el archivo original a la carpeta de destino
        shutil.copy2(archivo_original, archivo_destino)

        # Modificar el archivo en la carpeta de destino
        modificar_etiqueta(archivo_destino, nro_lote, fecha_vencimiento)

        return archivo_destino

    except Exception as e:
        raise Exception(f"Error al procesar la etiqueta: {str(e)}") 